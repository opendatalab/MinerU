# Copyright (c) Opendatalab. All rights reserved.
from io import BytesIO

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.annotations import Link

from mineru.backend.pipeline.pipeline_middle_json_mkcontent import union_make
from mineru.model.docx.docx_converter import DocxConverter
from mineru.model.pptx.pptx_converter import PptxConverter
from mineru.utils.enum_class import BlockType, ContentType, MakeMode
from mineru.utils.pdf_hyperlink import enrich_pdf_hyperlinks


def test_pptx_table_preserves_run_hyperlinks() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        1, 1, Inches(1), Inches(1), Inches(4), Inches(1)
    ).table
    paragraph = table.cell(0, 0).text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "Example & docs"
    run.hyperlink.address = "https://example.com/docs"

    stream = BytesIO()
    presentation.save(stream)
    stream.seek(0)
    converter = PptxConverter()
    converter.convert(stream)

    table_html = next(
        block["content"]
        for page in converter.pages
        for block in page
        if block["type"] == BlockType.TABLE
    )
    assert (
        '<a href="https://example.com/docs">Example &amp; docs</a>'
        in table_html
    )


def test_docx_native_table_fallback_preserves_hyperlinks() -> None:
    document = Document()
    table = document.add_table(rows=1, cols=1)
    paragraph = table.cell(0, 0).paragraphs[0]
    relationship_id = document.part.relate_to(
        "https://example.com/issues",
        RELATIONSHIP_TYPE.HYPERLINK,
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    text = OxmlElement("w:t")
    text.text = "Issue tracker"
    run.append(text)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)

    converter = DocxConverter()
    converter.docx_obj = document
    table_html = converter._build_table_html_from_xml(table._tbl)

    assert (
        '<a href="https://example.com/issues">Issue tracker</a>'
        in table_html
    )


def test_pdf_annotations_are_rendered_as_markdown_links_and_anchors() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_blank_page(width=200, height=200)
    writer.add_annotation(
        page_number=0,
        annotation=Link(
            rect=(10, 150, 90, 170),
            url="https://example.com/report",
        ),
    )
    writer.add_annotation(
        page_number=0,
        annotation=Link(
            rect=(100, 150, 190, 170),
            target_page_index=1,
        ),
    )
    stream = BytesIO()
    writer.write(stream)

    pdf_info = [
        {
            "page_idx": 0,
            "page_size": [200, 200],
            "para_blocks": [
                {
                    "type": BlockType.TEXT,
                    "index": 0,
                    "bbox": [10, 30, 190, 50],
                    "lines": [
                        {
                            "bbox": [10, 30, 190, 50],
                            "spans": [
                                {
                                    "type": ContentType.TEXT,
                                    "content": "External",
                                    "bbox": [10, 30, 90, 50],
                                },
                                {
                                    "type": ContentType.TEXT,
                                    "content": "Next page",
                                    "bbox": [100, 30, 190, 50],
                                },
                            ],
                        }
                    ],
                }
            ],
            "discarded_blocks": [],
        },
        {
            "page_idx": 1,
            "page_size": [200, 200],
            "para_blocks": [],
            "discarded_blocks": [],
        },
    ]

    enrich_pdf_hyperlinks(pdf_info, stream.getvalue())
    spans = pdf_info[0]["para_blocks"][0]["lines"][0]["spans"]

    assert spans[0]["type"] == ContentType.HYPERLINK
    assert spans[0]["url"] == "https://example.com/report"
    assert spans[1]["type"] == ContentType.HYPERLINK
    assert spans[1]["url"] == "#page-2"
    assert pdf_info[1]["page_anchor"] == "page-2"

    markdown = union_make(pdf_info, MakeMode.MM_MD)
    assert "[External](<https://example.com/report>)" in markdown
    assert "[Next page](<#page-2>)" in markdown
    assert '<a id="page-2"></a>' in markdown
