from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from mineru.model.flash import DocxModel, PptxModel, XlsxModel
from mineru.model.flash.pptx.pptx_converter import (
    PptxConverter,
    _PPTX_TITLE_ROLE_CENTER,
    _PPTX_TITLE_ROLE_KEY,
    _PPTX_TITLE_ROLE_SUBTITLE,
    _PPTX_TITLE_ROLE_TITLE,
)
from mineru.types import BlockType


_PROJECT_ROOT = Path(__file__).resolve().parents[2]
_OFFICE_SAMPLE_DIR = _PROJECT_ROOT / "demo" / "office_docs"


def _predict_sample(model: Any, suffix: str) -> list[list[dict[str, Any]]]:
    """调用指定 Office model 解析仓库中的真实样例。"""
    with (_OFFICE_SAMPLE_DIR / f"{suffix}_01.{suffix}").open("rb") as stream:
        return model.predict(stream)


def _flatten_pages(pages: list[list[dict[str, Any]]]) -> list[dict[str, Any]]:
    """按页面顺序展开 model_output blocks。"""
    return [block for page in pages for block in page]


def _find_block_by_content(blocks: list[dict[str, Any]], content: str) -> dict[str, Any]:
    """按完整 content 查找唯一 block。"""
    matches = [block for block in blocks if block.get("content") == content]
    assert len(matches) == 1
    return matches[0]


def _find_block_containing(blocks: list[dict[str, Any]], content: str) -> dict[str, Any]:
    """按 content 子串查找唯一 block。"""
    matches = [block for block in blocks if isinstance(block.get("content"), str) and content in block["content"]]
    assert len(matches) == 1
    return matches[0]


def _set_subtitle(slide: Any, text: str) -> None:
    """设置标题页的 Subtitle 占位符文本。"""
    subtitle = next(shape for shape in slide.placeholders if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE)
    subtitle.text = text


def test_docx_model_splits_document_and_paragraph_titles() -> None:
    """验证 DOCX model 按样式来源拆分文档标题和段落标题。"""
    pages = _predict_sample(DocxModel(), "docx")
    blocks = _flatten_pages(pages)

    assert all(block.get("type") != BlockType.TITLE for block in blocks)
    assert pages[0][0]["content"] == "MinerU supports DOCX document parsing now"

    doc_title = _find_block_by_content(blocks, "MinerU supports DOCX document parsing now")
    assert doc_title["type"] == BlockType.DOC_TITLE
    assert set(doc_title) <= {"type", "content", "anchor"}
    assert "level" not in doc_title
    assert "is_numbered_style" not in doc_title

    heading = _find_block_by_content(blocks, "后面是正常的章节标题")
    assert heading["type"] == BlockType.PARAGRAPH_TITLE
    assert heading["level"] == 1
    assert heading["is_numbered_style"] is False

    numbered_heading = _find_block_containing(blocks, "有序列表构成的章节标题")
    assert numbered_heading["type"] == BlockType.PARAGRAPH_TITLE
    assert numbered_heading["level"] == 1
    assert numbered_heading["is_numbered_style"] is True


def test_pptx_model_splits_title_placeholders_and_subtitle() -> None:
    """验证 PPTX 真实样例的首页主标题、Subtitle 和普通页标题映射。"""
    pages = _predict_sample(PptxModel(), "pptx")
    blocks = _flatten_pages(pages)

    assert all(block.get("type") != BlockType.TITLE for block in blocks)
    assert all(_PPTX_TITLE_ROLE_KEY not in block for block in blocks)
    assert [block["content"] for block in pages[0] if block.get("content") in {"Test Table Slide", "With footnote"}] == [
        "Test Table Slide",
        "With footnote",
    ]

    doc_title = _find_block_by_content(pages[0], "Test Table Slide")
    assert doc_title["type"] == BlockType.DOC_TITLE
    assert set(doc_title) <= {"type", "content", "anchor"}
    assert "level" not in doc_title
    assert "is_numbered_style" not in doc_title

    subtitle = _find_block_by_content(pages[0], "With footnote")
    assert subtitle["type"] == BlockType.TEXT
    assert "level" not in subtitle

    slide_title = _find_block_by_content(pages[1], "Second slide title")
    assert slide_title["type"] == BlockType.PARAGRAPH_TITLE
    assert slide_title["level"] == 2


def test_pptx_internal_titles_are_finalized_without_losing_levels() -> None:
    """验证 PPTX 内部 title 收口时清理文档标题字段并保留段落标题层级。"""
    blocks = [
        {
            "type": BlockType.TITLE,
            "content": "document",
            "level": 2,
            "is_numbered_style": False,
            _PPTX_TITLE_ROLE_KEY: _PPTX_TITLE_ROLE_CENTER,
        },
        {
            "type": BlockType.TITLE,
            "content": "subtitle",
            "level": 2,
            _PPTX_TITLE_ROLE_KEY: _PPTX_TITLE_ROLE_SUBTITLE,
        },
        {
            "type": BlockType.TITLE,
            "content": "slide title",
            "level": 2,
            _PPTX_TITLE_ROLE_KEY: _PPTX_TITLE_ROLE_TITLE,
        },
        {
            "type": BlockType.TITLE,
            "content": "promoted title",
            "level": 3,
        },
    ]

    PptxConverter._finalize_slide_title_types(blocks, is_first_visible_slide=True)

    assert [block["type"] for block in blocks] == [
        BlockType.DOC_TITLE,
        BlockType.TEXT,
        BlockType.PARAGRAPH_TITLE,
        BlockType.PARAGRAPH_TITLE,
    ]
    assert "level" not in blocks[0]
    assert "is_numbered_style" not in blocks[0]
    assert "level" not in blocks[1]
    assert blocks[2]["level"] == 2
    assert blocks[3]["level"] == 3
    assert all(_PPTX_TITLE_ROLE_KEY not in block for block in blocks)

    later_center_title = [
        {
            "type": BlockType.TITLE,
            "content": "later center title",
            "level": 2,
            _PPTX_TITLE_ROLE_KEY: _PPTX_TITLE_ROLE_CENTER,
        }
    ]
    PptxConverter._finalize_slide_title_types(
        later_center_title,
        is_first_visible_slide=False,
    )
    assert later_center_title[0]["type"] == BlockType.PARAGRAPH_TITLE
    assert later_center_title[0]["level"] == 2


def test_pptx_notes_only_slide_does_not_consume_document_title() -> None:
    """验证空白且仅含备注的前置页不会抢占首个可见文档标题。"""
    presentation = Presentation()
    notes_only_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    notes_only_slide.notes_slide.notes_text_frame.text = "Presenter note"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Document title"
    _set_subtitle(title_slide, "Document subtitle")

    later_title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    later_title_slide.shapes.title.text = "Section cover"

    stream = BytesIO()
    presentation.save(stream)
    stream.seek(0)
    pages = PptxModel().predict(stream)
    blocks = _flatten_pages(pages)

    assert _find_block_by_content(blocks, "Presenter note")["type"] == BlockType.PAGE_FOOTNOTE
    assert _find_block_by_content(blocks, "Document title")["type"] == BlockType.DOC_TITLE
    assert _find_block_by_content(blocks, "Document subtitle")["type"] == BlockType.TEXT
    later_title = _find_block_by_content(blocks, "Section cover")
    assert later_title["type"] == BlockType.PARAGRAPH_TITLE
    assert later_title["level"] == 2
    assert all(block.get("type") != BlockType.TITLE for block in blocks)


def test_xlsx_model_emits_sheet_names_as_paragraph_titles() -> None:
    """验证 XLSX 多 Sheet 名称全部输出为无 level 的段落标题。"""
    pages = _predict_sample(XlsxModel(), "xlsx")
    blocks = _flatten_pages(pages)
    sheet_titles = [_find_block_by_content(blocks, sheet_name) for sheet_name in ("Sheet1", "Sheet2", "Sheet3")]

    assert [page[0]["content"] for page in pages] == ["Sheet1", "Sheet2", "Sheet3"]
    assert all(block["type"] == BlockType.PARAGRAPH_TITLE for block in sheet_titles)
    assert all("level" not in block for block in sheet_titles)
    assert all(block.get("type") not in {BlockType.TITLE, BlockType.DOC_TITLE} for block in blocks)
