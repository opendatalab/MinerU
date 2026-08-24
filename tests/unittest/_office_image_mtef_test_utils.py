"""把带 MTEF comment 的图片注入确定性 DOCX/PPTX/XLSX。"""

from __future__ import annotations

import base64
from io import BytesIO
from zipfile import ZipFile

from docx import Document
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XlsxImage
from pptx import Presentation
from pptx.util import Inches
from lxml import etree  # type: ignore[reportAttributeAccessIssue]

from _ooxml_mtef_test_utils import (
    IMAGE_REL_TYPE,
    PACKAGE_RELS_NS,
    REL_NS,
    _rewrite_zip,
)

_VALID_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/"
    "x8AAusB9Wl2l9sAAAAASUVORK5CYII="
)


def _replace_media_parts(
    package: bytes,
    *,
    prefix: str,
    payload: bytes,
) -> bytes:
    """把指定 OOXML media 目录下的图片成员替换为原始测试载荷。"""

    with ZipFile(BytesIO(package)) as source:
        replacements = {
            name: payload
            for name in source.namelist()
            if name.startswith(prefix)
        }
    if not replacements:
        raise ValueError(f"OOXML image fixture has no media under {prefix}")
    return _rewrite_zip(package, replacements, {})


def build_image_docx(
    image_payload: bytes,
    *,
    inline: bool = False,
    table: bool = False,
    header: bool = False,
    paragraph_style: str | None = None,
) -> bytes:
    """构造正文、表格或页眉含单张图片的 DOCX。"""

    document = Document()
    if table:
        paragraph = document.add_table(rows=1, cols=1).cell(0, 0).paragraphs[0]
    elif header:
        paragraph = document.sections[0].header.paragraphs[0]
    else:
        paragraph = document.add_paragraph()
    if paragraph_style is not None:
        paragraph.style = paragraph_style
    if inline:
        paragraph.add_run("before ")
    paragraph.add_run().add_picture(BytesIO(_VALID_TINY_PNG))
    if inline:
        paragraph.add_run(" after")
    buffer = BytesIO()
    document.save(buffer)
    return _replace_media_parts(
        buffer.getvalue(),
        prefix="word/media/",
        payload=image_payload,
    )


def _move_pptx_picture_to_notes(package: bytes) -> bytes:
    """把单页 PPTX 的 picture shape 移到 notesSlide 并重绑图片关系。"""

    replacements: dict[str, bytes] = {}
    with ZipFile(BytesIO(package)) as source:
        slide_root = etree.fromstring(source.read("ppt/slides/slide1.xml"))
        notes_root = etree.fromstring(
            source.read("ppt/notesSlides/notesSlide1.xml")
        )
        pictures = slide_root.xpath(".//*[local-name()='pic']")
        note_trees = notes_root.xpath(".//*[local-name()='spTree']")
        if not pictures or not note_trees:
            raise ValueError("PPTX notes image fixture cannot locate picture trees")
        picture = pictures[0]
        parent = picture.getparent()
        parent.remove(picture)
        note_trees[0].append(picture)

        blips = picture.xpath(".//*[local-name()='blip']")
        if not blips:
            raise ValueError("PPTX notes image fixture has no blip")
        relationship_id = blips[0].get(f"{{{REL_NS}}}embed")
        slide_rels = etree.fromstring(
            source.read("ppt/slides/_rels/slide1.xml.rels")
        )
        target = next(
            (
                rel.get("Target")
                for rel in slide_rels
                if rel.get("Id") == relationship_id
            ),
            None,
        )
        if not target:
            raise ValueError("PPTX notes image relationship is missing")
        notes_rels_name = "ppt/notesSlides/_rels/notesSlide1.xml.rels"
        notes_rels = etree.fromstring(source.read(notes_rels_name))
        new_relationship_id = "rIdImageMtef"
        etree.SubElement(
            notes_rels,
            f"{{{PACKAGE_RELS_NS}}}Relationship",
            Id=new_relationship_id,
            Type=IMAGE_REL_TYPE,
            Target=target,
        )
        blips[0].set(f"{{{REL_NS}}}embed", new_relationship_id)
        replacements = {
            "ppt/slides/slide1.xml": etree.tostring(
                slide_root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            ),
            "ppt/notesSlides/notesSlide1.xml": etree.tostring(
                notes_root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            ),
            notes_rels_name: etree.tostring(
                notes_rels,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            ),
        }
    return _rewrite_zip(package, replacements, {})


def build_image_pptx(
    image_payload: bytes,
    *,
    notes: bool = False,
) -> bytes:
    """构造一页含单张图片 shape 或 notes 图片的 PPTX。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        BytesIO(_VALID_TINY_PNG),
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1.5),
    )
    if notes:
        _ = slide.notes_slide
    buffer = BytesIO()
    presentation.save(buffer)
    package = buffer.getvalue()
    if notes:
        package = _move_pptx_picture_to_notes(package)
    return _replace_media_parts(
        package,
        prefix="ppt/media/",
        payload=image_payload,
    )


def build_image_xlsx(
    image_payload: bytes,
    *,
    cell_value: str | None = None,
) -> bytes:
    """构造 A1 anchor 图片及可选单元格内容的 XLSX。"""

    workbook = Workbook()
    worksheet = workbook.active
    if worksheet is None:
        raise ValueError("XLSX image fixture has no active worksheet")
    worksheet.title = "Images"
    if cell_value is not None:
        worksheet["A1"] = cell_value
    image = XlsxImage(BytesIO(_VALID_TINY_PNG))
    worksheet.add_image(image, "A1")
    buffer = BytesIO()
    workbook.save(buffer)
    return _replace_media_parts(
        buffer.getvalue(),
        prefix="xl/media/",
        payload=image_payload,
    )
