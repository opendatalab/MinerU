"""构造带 Equation.3 OLE 对象的确定性 DOCX/PPTX/XLSX 测试包。"""

from __future__ import annotations

from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from docx import Document
from lxml import etree  # type: ignore[reportAttributeAccessIssue]
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from _mtef_test_utils import _TINY_PNG, build_equation_object

CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
PACKAGE_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
V_NS = "urn:schemas-microsoft-com:vml"
O_NS = "urn:schemas-microsoft-com:office:office"
X_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
XDR_NS = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
XVML_NS = "urn:schemas-microsoft-com:office:excel"
OLE_REL_TYPE = f"{REL_NS}/oleObject"
IMAGE_REL_TYPE = f"{REL_NS}/image"


def _rewrite_zip(
    package: bytes,
    replacements: dict[str, bytes],
    additions: dict[str, bytes],
) -> bytes:
    """在内存中替换或新增 OPC 成员并保持确定性压缩。"""

    source_buffer = BytesIO(package)
    output = BytesIO()
    with ZipFile(source_buffer) as source, ZipFile(output, "w", ZIP_DEFLATED) as target:
        existing = set()
        for info in source.infolist():
            existing.add(info.filename)
            payload = replacements.get(info.filename, source.read(info.filename))
            normalized = ZipInfo(info.filename, date_time=(1980, 1, 1, 0, 0, 0))
            normalized.compress_type = ZIP_DEFLATED
            normalized.external_attr = info.external_attr
            target.writestr(normalized, payload)
        for name, payload in additions.items():
            if name in existing:
                continue
            info = ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, payload)
    return output.getvalue()


def _content_types_with_ole(content_types: bytes) -> bytes:
    """为 .bin Equation.3 persistence part 添加默认内容类型。"""

    root = etree.fromstring(content_types)
    if not any(
        child.get("Extension", "").casefold() == "bin"
        for child in root
    ):
        etree.SubElement(
            root,
            f"{{{CONTENT_TYPES_NS}}}Default",
            Extension="bin",
            ContentType="application/vnd.openxmlformats-officedocument.oleObject",
        )
    if not any(
        child.get("Extension", "").casefold() == "png"
        for child in root
    ):
        etree.SubElement(
            root,
            f"{{{CONTENT_TYPES_NS}}}Default",
            Extension="png",
            ContentType="image/png",
        )
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _relationships_root(payload: bytes | None) -> etree._Element:
    """读取 relationships XML，不存在时创建空根。"""

    if payload:
        return etree.fromstring(payload)
    return etree.Element(f"{{{PACKAGE_RELS_NS}}}Relationships")


def _append_relationship(
    root: etree._Element,
    relationship_id: str,
    reltype: str,
    target: str,
    *,
    external: bool = False,
) -> None:
    """向 relationships 根追加一个确定性关系。"""

    attributes = {
        "Id": relationship_id,
        "Type": reltype,
        "Target": target,
    }
    if external:
        attributes["TargetMode"] = "External"
    etree.SubElement(
        root,
        f"{{{PACKAGE_RELS_NS}}}Relationship",
        **attributes,
    )


def _word_ole_run(
    object_index: int,
    ole_relationship_id: str,
    image_relationship_id: str,
) -> etree._Element:
    """构造包含 VML 预览和 o:OLEObject 的 Word run。"""

    run = etree.Element(f"{{{W_NS}}}r")
    object_element = etree.SubElement(run, f"{{{W_NS}}}object")
    shape_id = f"_x0000_i{1025 + object_index}"
    shape = etree.SubElement(
        object_element,
        f"{{{V_NS}}}shape",
        id=shape_id,
        style="width:120pt;height:36pt",
    )
    etree.SubElement(
        shape,
        f"{{{V_NS}}}imagedata",
        {f"{{{REL_NS}}}id": image_relationship_id},
    )
    etree.SubElement(
        object_element,
        f"{{{O_NS}}}OLEObject",
        {
            "Type": "Embed",
            "ProgID": "Equation.3",
            "ShapeID": shape_id,
            "DrawAspect": "Content",
            "ObjectID": f"_{1200000000 + object_index}",
            f"{{{REL_NS}}}id": ole_relationship_id,
        },
    )
    return run


def _replace_word_placeholder(
    root: etree._Element,
    placeholder: str,
    replacement: etree._Element,
) -> None:
    """用公式 run 或 AlternateContent 替换指定占位 run。"""

    text_node = next(
        (
            node
            for node in root.findall(f".//{{{W_NS}}}t")
            if node.text == placeholder
        ),
        None,
    )
    if text_node is None:
        raise ValueError(f"DOCX fixture placeholder is missing: {placeholder}")
    run = text_node.getparent()
    parent = run.getparent()
    parent.replace(run, replacement)


def _patch_word_part(
    source: ZipFile,
    part_name: str,
    formulas: list[tuple[str, bytes, bool]],
    additions: dict[str, bytes],
) -> tuple[bytes, bytes]:
    """向一个 Word XML part 及其 relationships 注入多个 Equation.3 对象。"""

    root = etree.fromstring(source.read(part_name))
    directory, basename = part_name.rsplit("/", 1)
    part_token = basename.rsplit(".", 1)[0]
    rels_name = f"{directory}/_rels/{basename}.rels"
    rels = _relationships_root(source.read(rels_name) if rels_name in source.namelist() else None)
    for index, (placeholder, mtef, alternate_omml) in enumerate(formulas, start=1):
        ole_rid = f"rIdMtefOle{index}"
        image_rid = f"rIdMtefImage{index}"
        object_part = f"word/embeddings/oleObjectMtef_{part_token}_{index}.bin"
        image_part = f"word/media/equationMtef_{part_token}_{index}.png"
        _append_relationship(
            rels,
            ole_rid,
            OLE_REL_TYPE,
            f"../embeddings/{object_part.rsplit('/', 1)[-1]}"
            if directory != "word"
            else f"embeddings/{object_part.rsplit('/', 1)[-1]}",
        )
        _append_relationship(
            rels,
            image_rid,
            IMAGE_REL_TYPE,
            f"../media/{image_part.rsplit('/', 1)[-1]}"
            if directory != "word"
            else f"media/{image_part.rsplit('/', 1)[-1]}",
        )
        additions[object_part] = build_equation_object(mtef)
        additions[image_part] = _TINY_PNG
        ole_run = _word_ole_run(index, ole_rid, image_rid)
        if alternate_omml:
            alternate = etree.Element(f"{{{MC_NS}}}AlternateContent")
            choice = etree.SubElement(
                alternate,
                f"{{{MC_NS}}}Choice",
                Requires="m",
            )
            math = etree.SubElement(choice, f"{{{M_NS}}}oMath")
            math_run = etree.SubElement(math, f"{{{M_NS}}}r")
            etree.SubElement(math_run, f"{{{M_NS}}}t").text = "z"
            fallback = etree.SubElement(alternate, f"{{{MC_NS}}}Fallback")
            fallback.append(ole_run)
            replacement = alternate
        else:
            replacement = ole_run
        _replace_word_placeholder(root, placeholder, replacement)
    return (
        etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True),
        etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True),
    )


def _wrap_word_equations_in_textboxes(document_xml: bytes) -> bytes:
    """把正文 Equation.3 run 包入 VML textbox，验证文本框独立遍历。"""

    root = etree.fromstring(document_xml)
    for index, ole_object in enumerate(
        root.findall(f".//{{{O_NS}}}OLEObject"),
        start=1,
    ):
        object_element = ole_object.getparent()
        equation_run = object_element.getparent()
        parent = equation_run.getparent()
        outer_run = etree.Element(f"{{{W_NS}}}r")
        pict = etree.SubElement(outer_run, f"{{{W_NS}}}pict")
        shape = etree.SubElement(
            pict,
            f"{{{V_NS}}}shape",
            id=f"_x0000_txbx{index}",
        )
        textbox = etree.SubElement(shape, f"{{{V_NS}}}textbox")
        content = etree.SubElement(textbox, f"{{{W_NS}}}txbxContent")
        paragraph = etree.SubElement(content, f"{{{W_NS}}}p")
        parent.replace(equation_run, outer_run)
        paragraph.append(equation_run)
    return etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )


def _set_word_equations_as_icons(document_xml: bytes) -> bytes:
    """把 Word OLE 对象切换为 DrawAspect=Icon。"""

    root = etree.fromstring(document_xml)
    for ole_object in root.findall(f".//{{{O_NS}}}OLEObject"):
        ole_object.set("DrawAspect", "Icon")
    return etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )


def build_equation_docx(
    formulas: list[bytes],
    *,
    inline: bool = False,
    table: bool = False,
    header_footer: bool = False,
    alternate_omml: bool = False,
    textbox: bool = False,
    show_as_icon: bool = False,
) -> bytes:
    """构造正文、表格或页眉页脚中含 Equation.3 对象的 DOCX。"""

    document = Document()
    target_part = "word/document.xml"
    placeholders: list[str] = []
    if table:
        cell = document.add_table(rows=1, cols=1).cell(0, 0)
        for index in range(len(formulas)):
            placeholder = f"[[MTEF_{index}]]"
            cell.paragraphs[0].add_run(placeholder)
            placeholders.append(placeholder)
    elif header_footer:
        header = document.sections[0].header
        footer = document.sections[0].footer
        for index in range(len(formulas)):
            placeholder = f"[[MTEF_{index}]]"
            target = header if index % 2 == 0 else footer
            target.paragraphs[0].add_run(placeholder)
            placeholders.append(placeholder)
        target_part = "word/header1.xml"
    else:
        for index in range(len(formulas)):
            paragraph = document.add_paragraph()
            if inline:
                paragraph.add_run("before ")
            placeholder = f"[[MTEF_{index}]]"
            paragraph.add_run(placeholder)
            placeholders.append(placeholder)
            if inline:
                paragraph.add_run(" after")

    buffer = BytesIO()
    document.save(buffer)
    package = buffer.getvalue()
    replacements: dict[str, bytes] = {}
    additions: dict[str, bytes] = {}
    with ZipFile(BytesIO(package)) as source:
        if header_footer:
            header_parts = sorted(
                name
                for name in source.namelist()
                if name.startswith("word/header") and name.endswith(".xml")
            )
            footer_parts = sorted(
                name
                for name in source.namelist()
                if name.startswith("word/footer") and name.endswith(".xml")
            )
            assignments: dict[str, list[tuple[str, bytes, bool]]] = {}
            for index, (placeholder, mtef) in enumerate(zip(placeholders, formulas, strict=True)):
                part = header_parts[0] if index % 2 == 0 else footer_parts[0]
                assignments.setdefault(part, []).append((placeholder, mtef, False))
            for part, entries in assignments.items():
                xml, rels = _patch_word_part(source, part, entries, additions)
                replacements[part] = xml
                directory, basename = part.rsplit("/", 1)
                rels_name = f"{directory}/_rels/{basename}.rels"
                if rels_name in source.namelist():
                    replacements[rels_name] = rels
                else:
                    additions[rels_name] = rels
        else:
            entries = [
                (
                    placeholder,
                    mtef,
                    alternate_omml and index == 0,
                )
                for index, (placeholder, mtef) in enumerate(
                    zip(placeholders, formulas, strict=True)
                )
            ]
            xml, rels = _patch_word_part(source, target_part, entries, additions)
            if textbox:
                xml = _wrap_word_equations_in_textboxes(xml)
            if show_as_icon:
                xml = _set_word_equations_as_icons(xml)
            replacements[target_part] = xml
            replacements["word/_rels/document.xml.rels"] = rels
        replacements["[Content_Types].xml"] = _content_types_with_ole(
            source.read("[Content_Types].xml")
        )
    return _rewrite_zip(package, replacements, additions)


def build_equation_pptx(
    formulas: list[bytes],
    *,
    show_as_icon: bool = False,
    notes: bool = False,
    alternate_omml: bool = False,
) -> bytes:
    """使用 python-pptx 生成每页一个 Equation.3 OLE 对象的 PPTX。"""

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, mtef in enumerate(formulas):
        slide = (
            presentation.slides[0]
            if index == 0
            else presentation.slides.add_slide(presentation.slide_layouts[6])
        )
        slide.shapes.add_ole_object(
            BytesIO(build_equation_object(mtef)),
            "Equation.3",
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1.5),
            icon_file=BytesIO(_TINY_PNG),
        )
        if notes:
            _ = slide.notes_slide
    buffer = BytesIO()
    presentation.save(buffer)
    package = buffer.getvalue()
    if notes:
        return _move_pptx_equations_to_notes(
            package,
            show_as_icon=show_as_icon,
        )
    if show_as_icon:
        return package

    replacements: dict[str, bytes] = {}
    with ZipFile(BytesIO(package)) as source:
        for name in source.namelist():
            if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                continue
            root = etree.fromstring(source.read(name))
            changed = False
            for ole_object in root.xpath(".//*[local-name()='oleObj']"):
                ole_object.set("showAsIcon", "0")
                changed = True
                if alternate_omml:
                    graphic_data = ole_object.getparent()
                    math = etree.SubElement(graphic_data, f"{{{M_NS}}}oMath")
                    math_run = etree.SubElement(math, f"{{{M_NS}}}r")
                    etree.SubElement(math_run, f"{{{M_NS}}}t").text = "z"
            if changed:
                replacements[name] = etree.tostring(
                    root,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone=True,
                )
    return _rewrite_zip(package, replacements, {})


def _move_pptx_equations_to_notes(
    package: bytes,
    *,
    show_as_icon: bool,
) -> bytes:
    """把每页 OLE graphicFrame 移到对应 notesSlide 并重建关系。"""

    replacements: dict[str, bytes] = {}
    with ZipFile(BytesIO(package)) as source:
        slide_parts = sorted(
            name
            for name in source.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        for index, slide_part in enumerate(slide_parts, start=1):
            notes_part = f"ppt/notesSlides/notesSlide{index}.xml"
            if notes_part not in source.namelist():
                continue
            slide_root = etree.fromstring(source.read(slide_part))
            graphic_frame = next(
                iter(slide_root.xpath(".//*[local-name()='graphicFrame'][.//*[local-name()='oleObj']]")),
                None,
            )
            if graphic_frame is None:
                continue
            for ole_object in graphic_frame.xpath(".//*[local-name()='oleObj']"):
                ole_object.set("showAsIcon", "1" if show_as_icon else "0")

            slide_rels_name = f"ppt/slides/_rels/slide{index}.xml.rels"
            slide_rels = _relationships_root(source.read(slide_rels_name))
            referenced_ids = {
                value
                for node in graphic_frame.iter()
                for attribute, value in node.attrib.items()
                if attribute in {f"{{{REL_NS}}}id", f"{{{REL_NS}}}embed"}
            }
            copied_relationships = [
                relationship
                for relationship in slide_rels
                if relationship.get("Id") in referenced_ids
            ]
            notes_rels_name = f"ppt/notesSlides/_rels/notesSlide{index}.xml.rels"
            notes_rels = _relationships_root(source.read(notes_rels_name))
            id_mapping: dict[str, str] = {}
            for relation_index, relationship in enumerate(copied_relationships, start=1):
                old_id = relationship.get("Id") or ""
                new_id = f"rIdMtef{relation_index}"
                id_mapping[old_id] = new_id
                _append_relationship(
                    notes_rels,
                    new_id,
                    relationship.get("Type") or "",
                    relationship.get("Target") or "",
                )
            for node in graphic_frame.iter():
                for attribute in (f"{{{REL_NS}}}id", f"{{{REL_NS}}}embed"):
                    old_id = node.get(attribute)
                    if old_id in id_mapping:
                        node.set(attribute, id_mapping[old_id])

            graphic_frame.getparent().remove(graphic_frame)
            notes_root = etree.fromstring(source.read(notes_part))
            shape_tree = next(
                iter(notes_root.xpath(".//*[local-name()='spTree']")),
                None,
            )
            if shape_tree is None:
                continue
            shape_tree.append(graphic_frame)
            replacements[slide_part] = etree.tostring(
                slide_root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            replacements[notes_part] = etree.tostring(
                notes_root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            replacements[notes_rels_name] = etree.tostring(
                notes_rels,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
    return _rewrite_zip(package, replacements, {})


def _xlsx_anchor(
    ole_object: etree._Element,
    row: int,
    col: int,
    preview_relationship_id: str,
) -> None:
    """向 x:oleObject 写入 objectPr/anchor/from/to。"""

    object_properties = etree.SubElement(
        ole_object,
        f"{{{X_NS}}}objectPr",
        {f"{{{REL_NS}}}id": preview_relationship_id},
    )
    anchor = etree.SubElement(object_properties, f"{{{X_NS}}}anchor")
    from_marker = etree.SubElement(anchor, f"{{{X_NS}}}from")
    etree.SubElement(from_marker, f"{{{X_NS}}}col").text = str(col)
    etree.SubElement(from_marker, f"{{{X_NS}}}row").text = str(row)
    to_marker = etree.SubElement(anchor, f"{{{X_NS}}}to")
    etree.SubElement(to_marker, f"{{{X_NS}}}col").text = str(col + 2)
    etree.SubElement(to_marker, f"{{{X_NS}}}row").text = str(row + 3)


def _xlsx_drawing_parts(
    shape_id: int,
    row: int,
    col: int,
    *,
    omml: bool = False,
) -> tuple[bytes, bytes]:
    """构造按 cNvPr id 绑定预览的 DrawingML part 及 relationships。"""

    root = etree.Element(f"{{{XDR_NS}}}wsDr", nsmap={"xdr": XDR_NS, "a": A_NS, "r": REL_NS})
    anchor = etree.SubElement(root, f"{{{XDR_NS}}}twoCellAnchor")
    from_marker = etree.SubElement(anchor, f"{{{XDR_NS}}}from")
    etree.SubElement(from_marker, f"{{{XDR_NS}}}col").text = str(col)
    etree.SubElement(from_marker, f"{{{XDR_NS}}}row").text = str(row)
    picture = etree.SubElement(anchor, f"{{{XDR_NS}}}pic")
    non_visual = etree.SubElement(picture, f"{{{XDR_NS}}}nvPicPr")
    etree.SubElement(
        non_visual,
        f"{{{XDR_NS}}}cNvPr",
        id=str(shape_id),
        name="Equation preview",
    )
    fill = etree.SubElement(picture, f"{{{XDR_NS}}}blipFill")
    etree.SubElement(fill, f"{{{A_NS}}}blip", {f"{{{REL_NS}}}embed": "rIdPreview"})
    if omml:
        math = etree.SubElement(anchor, f"{{{M_NS}}}oMath")
        math_run = etree.SubElement(math, f"{{{M_NS}}}r")
        etree.SubElement(math_run, f"{{{M_NS}}}t").text = "z"
    rels = _relationships_root(None)
    _append_relationship(rels, "rIdPreview", IMAGE_REL_TYPE, "../media/equationMtef.png")
    return (
        etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True),
        etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True),
    )


def _xlsx_vml_parts(
    shape_id: int,
    row: int,
    col: int,
) -> tuple[bytes, bytes]:
    """构造按 `_x0000_sNNN` 绑定预览的 VML drawing 与 relationships。"""

    root = etree.Element("xml", nsmap={"v": V_NS, "x": XVML_NS, "r": REL_NS})
    shape = etree.SubElement(root, f"{{{V_NS}}}shape", id=f"_x0000_s{shape_id}")
    etree.SubElement(shape, f"{{{V_NS}}}imagedata", {f"{{{REL_NS}}}id": "rIdPreview"})
    client_data = etree.SubElement(shape, f"{{{XVML_NS}}}ClientData")
    etree.SubElement(client_data, f"{{{XVML_NS}}}Anchor").text = (
        f"{col}, 0, {row}, 0, {col + 2}, 0, {row + 3}, 0"
    )
    rels = _relationships_root(None)
    _append_relationship(rels, "rIdPreview", IMAGE_REL_TYPE, "../media/equationMtef.png")
    return (
        etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True),
        etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True),
    )


def build_equation_xlsx(
    formulas: list[bytes],
    *,
    anchor_mode: str = "objectPr",
    cell_value: str | None = None,
    hidden: bool = False,
    linked: bool = False,
    show_as_icon: bool = False,
    alternate_omml: bool = False,
) -> bytes:
    """构造使用 objectPr、DrawingML 或 VML anchor 的 Equation.3 XLSX。"""

    workbook = Workbook()
    worksheet = workbook.active
    if worksheet is None:
        raise ValueError("XLSX fixture workbook has no active sheet")
    worksheet.title = "Equations"
    if cell_value is not None:
        worksheet.cell(row=1, column=1, value=cell_value)
    if hidden:
        worksheet.sheet_state = "hidden"
        workbook.create_sheet("Visible")
    buffer = BytesIO()
    workbook.save(buffer)
    package = buffer.getvalue()

    replacements: dict[str, bytes] = {}
    additions: dict[str, bytes] = {}
    with ZipFile(BytesIO(package)) as source:
        worksheet_part = "xl/worksheets/sheet1.xml"
        root = etree.fromstring(source.read(worksheet_part))
        ole_objects = etree.SubElement(root, f"{{{X_NS}}}oleObjects")
        rels_name = "xl/worksheets/_rels/sheet1.xml.rels"
        rels = _relationships_root(source.read(rels_name) if rels_name in source.namelist() else None)
        for index, mtef in enumerate(formulas, start=1):
            shape_id = 1024 + index
            ole_rid = f"rIdMtefOle{index}"
            preview_rid = f"rIdMtefPreview{index}"
            ole_object = etree.SubElement(
                ole_objects,
                f"{{{X_NS}}}oleObject",
                {
                    "progId": "Equation.3",
                    "shapeId": str(shape_id),
                    f"{{{REL_NS}}}id": ole_rid,
                },
            )
            if show_as_icon:
                ole_object.set("dvAspect", "DVASPECT_ICON")
            if anchor_mode == "objectPr":
                _xlsx_anchor(ole_object, index - 1, 0, preview_rid)
                _append_relationship(
                    rels,
                    preview_rid,
                    IMAGE_REL_TYPE,
                    "../media/equationMtef.png",
                )
            _append_relationship(
                rels,
                ole_rid,
                OLE_REL_TYPE,
                (
                    "https://example.test/equation.bin"
                    if linked
                    else f"../embeddings/oleObjectMtef{index}.bin"
                ),
                external=linked,
            )
            if linked:
                ole_object.set("link", "https://example.test/equation.bin")
            else:
                additions[f"xl/embeddings/oleObjectMtef{index}.bin"] = build_equation_object(mtef)

        if anchor_mode == "drawing":
            etree.SubElement(root, f"{{{X_NS}}}drawing", {f"{{{REL_NS}}}id": "rIdMtefDrawing"})
            _append_relationship(
                rels,
                "rIdMtefDrawing",
                f"{REL_NS}/drawing",
                "../drawings/drawingMtef.xml",
            )
            drawing, drawing_rels = _xlsx_drawing_parts(
                1025,
                0,
                0,
                omml=alternate_omml,
            )
            additions["xl/drawings/drawingMtef.xml"] = drawing
            additions["xl/drawings/_rels/drawingMtef.xml.rels"] = drawing_rels
        elif anchor_mode == "vml":
            etree.SubElement(root, f"{{{X_NS}}}legacyDrawing", {f"{{{REL_NS}}}id": "rIdMtefVml"})
            _append_relationship(
                rels,
                "rIdMtefVml",
                f"{REL_NS}/vmlDrawing",
                "../drawings/vmlDrawingMtef.vml",
            )
            drawing, drawing_rels = _xlsx_vml_parts(1025, 0, 0)
            additions["xl/drawings/vmlDrawingMtef.vml"] = drawing
            additions["xl/drawings/_rels/vmlDrawingMtef.vml.rels"] = drawing_rels
        elif anchor_mode not in {"objectPr", "none"}:
            raise ValueError(f"unsupported XLSX equation anchor fixture: {anchor_mode}")

        additions["xl/media/equationMtef.png"] = _TINY_PNG
        replacements[worksheet_part] = etree.tostring(
            root,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )
        rels_payload = etree.tostring(
            rels,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )
        if rels_name in source.namelist():
            replacements[rels_name] = rels_payload
        else:
            additions[rels_name] = rels_payload
        replacements["[Content_Types].xml"] = _content_types_with_ole(
            source.read("[Content_Types].xml")
        )
    return _rewrite_zip(package, replacements, additions)
